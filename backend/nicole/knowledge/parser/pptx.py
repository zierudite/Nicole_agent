"""PPTX 解析器。

使用 python-pptx 解析 PowerPoint 文件，提取幻灯片文本和备注。
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from .base import BaseParser, ParserResult

logger = logging.getLogger(__name__)


class PPTXParser(BaseParser):
    """PowerPoint 解析器 (python-pptx)。"""

    def __init__(
        self,
        use_paddlex: bool = False,
        use_rapidocr: bool = False,
    ):
        self.use_paddlex = use_paddlex
        self.use_rapidocr = use_rapidocr

    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    async def parse(self, file_path: str) -> Optional[ParserResult]:
        """解析 .pptx 文件。"""
        if not self.validate_file(file_path):
            return None

        path = Path(file_path)
        metadata = {"filename": path.name, "size": path.stat().st_size}

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            slides_text = []
            pages = []

            for i, slide in enumerate(prs.slides):
                slide_content = []
                layout_name = slide.slide_layout.name

                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_content.append(text)

                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            slide_content.append(" | ".join(cells))

                slide_text = "\n".join(slide_content)
                slides_text.append(
                    f"--- Slide {i + 1} (Layout: {layout_name}) ---\n{slide_text}"
                )
                pages.append({
                    "page": i + 1,
                    "layout": layout_name,
                    "text": slide_text,
                })

            # 提取备注
            notes = []
            for i, slide in enumerate(prs.slides):
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes.append(f"--- Slide {i + 1} Notes ---\n{notes_text}")

            full_text = "\n\n".join(slides_text)
            if notes:
                full_text += "\n\n" + "\n\n".join(notes)

            return ParserResult(
                text=full_text,
                file_type="pptx",
                metadata={
                    **metadata,
                    "slides": len(slides_text),
                    "engine": "python-pptx",
                },
                pages=pages,
            )

        except ImportError:
            logger.error("python-pptx not installed")
            return None
        except Exception as e:
            logger.error(f"PPTX parse failed: {e}")
            return None
